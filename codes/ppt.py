from pptx import Presentation

def read_ppt(file):
    """Read a PPTX file and return its content as a string."""
    presentation = Presentation(file)
    document_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                document_text += shape.text + "\n"
    return document_text
